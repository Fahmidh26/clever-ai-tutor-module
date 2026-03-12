from __future__ import annotations

import re
from pathlib import Path

import httpx

from app.config import settings


class KBPipelineError(Exception):
    pass


def extract_text_from_document(file_path: Path, file_type: str) -> str:
    ext = (file_type or "").strip().lower()

    if ext in {"txt", "md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    if ext == "pdf":
        try:
            from pypdf import PdfReader  # type: ignore
        except Exception as exc:
            raise KBPipelineError("PDF extraction requires 'pypdf' dependency") from exc
        reader = PdfReader(str(file_path))
        parts: list[str] = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n".join(parts)

    if ext == "docx":
        try:
            from docx import Document  # type: ignore
        except Exception as exc:
            raise KBPipelineError("DOCX extraction requires 'python-docx' dependency") from exc
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    if ext == "pptx":
        try:
            from pptx import Presentation  # type: ignore
        except Exception as exc:
            raise KBPipelineError("PPTX extraction requires 'python-pptx' dependency") from exc
        prs = Presentation(str(file_path))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    lines.append(text)
        return "\n".join(lines)

    raise KBPipelineError(f"Unsupported document type: {ext}")


def clean_text(raw_text: str) -> str:
    text = (raw_text or "").replace("\x00", " ")
    text = re.sub(r"\r\n?", "\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, *, chunk_size: int, chunk_overlap: int, max_chunks: int) -> list[str]:
    source = (text or "").strip()
    if not source:
        return []

    if chunk_overlap >= chunk_size:
        chunk_overlap = max(0, chunk_size // 5)
    step = max(1, chunk_size - chunk_overlap)

    chunks: list[str] = []
    i = 0
    while i < len(source) and len(chunks) < max_chunks:
        segment = source[i : i + chunk_size].strip()
        if segment:
            chunks.append(segment)
        i += step
    return chunks


async def embed_chunks(chunks: list[str]) -> list[list[float]]:
    if not settings.openai_api_key:
        raise KBPipelineError("OPENAI_API_KEY is required for embeddings")
    if not chunks:
        return []

    payload = {
        "model": settings.embedding_model,
        "input": chunks,
    }
    headers = {
        "Authorization": f"Bearer {settings.openai_api_key}",
        "Content-Type": "application/json",
    }
    async with httpx.AsyncClient(timeout=settings.embedding_timeout_seconds) as client:
        response = await client.post(
            f"{settings.openai_base_url.rstrip('/')}/embeddings",
            headers=headers,
            json=payload,
        )
    if response.status_code >= 400:
        raise KBPipelineError(f"Embedding API error {response.status_code}: {response.text[:300]}")

    data = response.json().get("data", [])
    vectors = [item.get("embedding") for item in data if isinstance(item, dict)]
    if len(vectors) != len(chunks):
        raise KBPipelineError("Embedding response length mismatch")
    return vectors

